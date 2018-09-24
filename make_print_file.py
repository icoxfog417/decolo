import urllib.parse
import tempfile
import csv
import qrcode
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


host = "https://decolo.herokuapp.com"


def make_file(path, urls, names):
    pres = Presentation()

    for url, name in zip(urls, names):
        add_slide(pres, url, name)

    pres.save(path)


def add_slide(pres, url, name):
    layout = pres.slide_layouts[1]
    slide = pres.slides.add_slide(layout)

    shapes = slide.shapes
    title_shape = shapes.title
    title_shape.text = "QRコードからサイトにアクセスしてください\nアクセス後、カメラをオンにして顔を映して\nお楽しみください"
    title_shape.text_frame.paragraphs[0].font.size = Pt(30)
    title_shape.text_frame.paragraphs[0].font.name = "メイリオ"

    with tempfile.NamedTemporaryFile(delete=False) as tmpf:
        image = qrcode.make(url)
        image.save(tmpf)
        pic = slide.shapes.add_picture(tmpf.name, Cm(7.7), Cm(5.39), Cm(10), Cm(10))

    nameBox = slide.shapes.add_textbox(Cm(16.38), Cm(16), Cm(7.75), Cm(1.03))
    tf = nameBox.text_frame
    p = tf.add_paragraph()
    p.text = name + "様"
    p.font.name = "メイリオ"
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.RIGHT


if __name__ == "__main__":
    def make_urls():
        names = []
        urls = []
        with open("names.csv", encoding="utf-8") as f:
            reader = csv.reader(f, delimiter="\t")
            next(reader)
            for row in reader:
                name, disp_name, gender = row
                url = host + "?" + "&name={}&gender={}".format(
                        urllib.parse.quote(disp_name), gender)
                urls.append(url)
                names.append(name)
        return names, urls

    names, urls = make_urls()
    make_file("print.pptx", urls, names)
